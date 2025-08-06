# outputs/ppt_generator.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os, tempfile, datetime

def generate_ppt(df):
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 ratio
    prs.slide_height = Inches(7.5)

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Amazon SP N-Gram Insights"
        slide.placeholders[1].text = "Automated Ad Term Analysis\nGenerated on " + datetime.date.today().strftime("%b %d, %Y")

    def add_tool_summary_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "What This Tool Does"
        content = slide.placeholders[1].text_frame
        content.clear()

        content.add_paragraph().text = "🔹 Input: Amazon Search Term Impression Share report"
        content.add_paragraph().text = "🧠 Processing: N-gram tokenization → Context tagging → Metric calculations"
        content.add_paragraph().text = "📦 Output: Excel files, Bubble charts, PowerPoint insights"

    def add_summary_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Strategic Summary"
        content = slide.placeholders[1].text_frame
        content.clear()

        content.add_paragraph().text = "🔺 Scale: Terms with CE ≥ 110% of break-even"
        content.add_paragraph().text = "⚠️ Trim: Terms with CE ≤ 90% of break-even"
        content.add_paragraph().text = "🧮 Zones: High-efficiency, Break-even (±10%), Underperforming"

    def add_ngram_insight_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "n-Gram Insights"
        top_margin = Inches(1.5)
        col_width = prs.slide_width / 3

        for i, n in enumerate([1, 2, 3]):
            sub = df[df['n'] == n]
            top = sub[sub['Efficiency Zone'] == 'High-efficiency / Low-cost'].nlargest(3, 'CE')
            low = sub[sub['Efficiency Zone'] == 'Overpriced / Underperforming'].nsmallest(3, 'CE')

            left = Inches(i * (13.33 / 3))
            box = slide.shapes.add_textbox(left, top_margin, col_width, Inches(5))
            tf = box.text_frame
            tf.word_wrap = True

            p = tf.add_paragraph()
            p.text = f"{n}-Gram"
            p.font.size = Pt(14)
            p.font.bold = True

            tf.add_paragraph().text = "Top Efficient: " + ", ".join(top['ngram'])
            tf.add_paragraph().text = "Underperformers: " + ", ".join(low['ngram'])
            tf.add_paragraph().text = "Rec: Scale top CE terms; review or trim low CE terms."

    def add_cross_ngram_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Cross-n-Gram Patterns"
        content = slide.placeholders[1].text_frame
        content.clear()

        all_under = df[df['Efficiency Zone'] == 'Overpriced / Underperforming']['ngram']
        common_lows = all_under.value_counts()[all_under.value_counts() > 1].index.tolist()[:5]

        all_high = df[df['Efficiency Zone'] == 'High-efficiency / Low-cost']['ngram']
        common_highs = all_high.value_counts()[all_high.value_counts() > 1].index.tolist()[:5]

        content.add_paragraph().text = "⚠ Shared Underperformers: " + ", ".join(common_lows)
        content.add_paragraph().text = "💡 Repeated High Performers: " + ", ".join(common_highs)
        content.add_paragraph().text = "📌 Action: Pause shared low CE terms; invest in repeat winners."

    def add_appendix_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Appendix"

    # Slide sequence
    add_title_slide()
    add_tool_summary_slide()
    add_summary_slide()
    add_ngram_insight_slide()
    add_cross_ngram_slide()
    add_appendix_slide()

    path = os.path.join(tempfile.gettempdir(), "ngram_insights.pptx")
    prs.save(path)
    return path
